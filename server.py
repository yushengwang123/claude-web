import asyncio
import ipaddress
import json
import os
import re
import socket
import sqlite3
import time
import urllib.error
import urllib.request
import uuid
from collections import defaultdict
from contextlib import asynccontextmanager, contextmanager
from html.parser import HTMLParser
from pathlib import Path
from typing import Dict, Iterator, List, Optional, Set
from urllib.parse import urlparse

from fastapi import FastAPI, File, HTTPException, Query, UploadFile
from fastapi.responses import FileResponse, Response, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

from claude_web import __version__

_PKG_DIR = Path(__file__).parent
_DATA_DIR = Path(os.environ.get("CLAUDE_WEB_DATA_DIR", "")).resolve() if os.environ.get("CLAUDE_WEB_DATA_DIR") else Path.cwd()

STATIC_DIR = _PKG_DIR / "static"
HISTORY_DIR = _DATA_DIR / "history"
UPLOADS_DIR = _DATA_DIR / "uploads"
DB_PATH = _DATA_DIR / "claude-web.db"

HISTORY_DIR.mkdir(parents=True, exist_ok=True)
UPLOADS_DIR.mkdir(parents=True, exist_ok=True)

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
MAX_UPLOAD_MB = 20
IGNORED_DIRS = {".git", "node_modules", ".venv", "venv", "__pycache__", ".next", "dist", "build", ".cache", ".idea", ".vscode"}
KNOWN_TOOL_NAMES = {
    "Bash", "Read", "Write", "Edit", "MultiEdit", "Grep", "Glob",
    "WebFetch", "WebSearch", "TodoWrite", "Task", "NotebookEdit",
}

_running_processes: Dict[str, asyncio.subprocess.Process] = {}
_stopped_sessions: Set[str] = set()


async def _terminate_process(process: asyncio.subprocess.Process, grace: float = 3.0) -> None:
    if process.returncode is not None:
        return
    try:
        process.terminate()
    except ProcessLookupError:
        return
    try:
        await asyncio.wait_for(process.wait(), timeout=grace)
        return
    except asyncio.TimeoutError:
        pass
    try:
        process.kill()
    except ProcessLookupError:
        return
    try:
        await asyncio.wait_for(process.wait(), timeout=2.0)
    except asyncio.TimeoutError:
        pass


async def _shutdown_terminate_running_processes() -> None:
    if not _running_processes:
        return
    processes = list(_running_processes.values())
    _running_processes.clear()
    await asyncio.gather(
        *(_terminate_process(p) for p in processes),
        return_exceptions=True,
    )


_UPLOAD_RETENTION_SECONDS = 30 * 24 * 60 * 60  # 30 days


def _prune_old_uploads(retention_seconds: int = _UPLOAD_RETENTION_SECONDS) -> int:
    """Delete files in UPLOADS_DIR older than retention_seconds. Returns count
    of files removed. Best-effort: silently skips entries we can't stat/unlink."""
    if not UPLOADS_DIR.exists():
        return 0
    cutoff = time.time() - retention_seconds
    removed = 0
    for entry in UPLOADS_DIR.iterdir():
        if not entry.is_file():
            continue
        try:
            if entry.stat().st_mtime < cutoff:
                entry.unlink()
                removed += 1
        except OSError:
            continue
    return removed


@asynccontextmanager
async def _lifespan(app: FastAPI):
    # Prune stale uploads in a background thread so startup isn't blocked on disk IO.
    asyncio.get_event_loop().run_in_executor(None, _prune_old_uploads)
    try:
        yield
    finally:
        await _shutdown_terminate_running_processes()


app = FastAPI(title="Claude Code Web", lifespan=_lifespan)


async def _drain_stream(stream: asyncio.StreamReader, buffer: bytearray, limit: int = 256 * 1024) -> None:
    try:
        while True:
            chunk = await stream.read(8192)
            if not chunk:
                return
            remaining = limit - len(buffer)
            if remaining > 0:
                buffer.extend(chunk[:remaining])
    except asyncio.CancelledError:
        raise
    except Exception:
        return


_DB_INITIALIZED = False


@contextmanager
def db_connect() -> Iterator[sqlite3.Connection]:
    global _DB_INITIALIZED
    conn = sqlite3.connect(DB_PATH, timeout=10)
    try:
        conn.row_factory = sqlite3.Row
        if not _DB_INITIALIZED:
            conn.execute("PRAGMA journal_mode=WAL")
            conn.execute("PRAGMA synchronous=NORMAL")
            _DB_INITIALIZED = True
        conn.execute("PRAGMA busy_timeout=5000")
        yield conn
        conn.commit()
    except Exception:
        conn.rollback()
        raise
    finally:
        conn.close()


def ensure_column(conn: sqlite3.Connection, table: str, column: str, ddl: str) -> None:
    cols = {row["name"] for row in conn.execute(f"PRAGMA table_info({table})").fetchall()}
    if column not in cols:
        conn.execute(f"ALTER TABLE {table} ADD COLUMN {column} {ddl}")


def init_db() -> None:
    with db_connect() as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS sessions (
                id TEXT PRIMARY KEY,
                title TEXT NOT NULL DEFAULT '',
                cwd TEXT NOT NULL DEFAULT '',
                created_at REAL NOT NULL,
                updated_at REAL NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS prompts (
                id TEXT PRIMARY KEY,
                name TEXT NOT NULL,
                content TEXT NOT NULL,
                created_at REAL NOT NULL
            )
            """
        )
        ensure_column(conn, "sessions", "pinned", "INTEGER NOT NULL DEFAULT 0")
        ensure_column(conn, "sessions", "archived", "INTEGER NOT NULL DEFAULT 0")
        ensure_column(conn, "sessions", "tags", "TEXT NOT NULL DEFAULT ''")
        ensure_column(conn, "sessions", "manual_title", "INTEGER NOT NULL DEFAULT 0")
        ensure_column(conn, "sessions", "remote_session_id", "TEXT NOT NULL DEFAULT ''")
        ensure_column(conn, "sessions", "remote_ready", "INTEGER NOT NULL DEFAULT 0")


init_db()


def upsert_session(session_id: str, title: str, cwd: str) -> None:
    now = time.time()
    with db_connect() as conn:
        row = conn.execute(
            "SELECT title, manual_title FROM sessions WHERE id = ?", (session_id,)
        ).fetchone()
        if row is None:
            conn.execute(
                "INSERT INTO sessions (id, title, cwd, created_at, updated_at) VALUES (?, ?, ?, ?, ?)",
                (session_id, title, cwd, now, now),
            )
        else:
            new_title = row["title"]
            if not row["manual_title"] and not new_title:
                new_title = title
            conn.execute(
                "UPDATE sessions SET title = ?, cwd = ?, updated_at = ? WHERE id = ?",
                (new_title, cwd, now, session_id),
            )


def append_event(session_id: str, event: dict) -> None:
    path = HISTORY_DIR / f"{session_id}.jsonl"
    with path.open("a", encoding="utf-8") as f:
        f.write(json.dumps(event, ensure_ascii=False) + "\n")


def load_events(session_id: str) -> List[dict]:
    path = HISTORY_DIR / f"{session_id}.jsonl"
    if not path.exists():
        return []
    events: List[dict] = []
    with path.open("r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            try:
                events.append(json.loads(line))
            except json.JSONDecodeError:
                continue
    return events


def save_events(session_id: str, events: List[dict]) -> None:
    path = HISTORY_DIR / f"{session_id}.jsonl"
    if not events:
        if path.exists():
            path.unlink()
        return
    tmp_path = path.with_suffix(path.suffix + f".tmp.{os.getpid()}.{uuid.uuid4().hex[:8]}")
    try:
        with tmp_path.open("w", encoding="utf-8") as f:
            for event in events:
                f.write(json.dumps(event, ensure_ascii=False) + "\n")
            f.flush()
            os.fsync(f.fileno())
        os.replace(tmp_path, path)
    except Exception:
        tmp_path.unlink(missing_ok=True)
        raise


def summarize_text_from_events(events: List[dict]) -> str:
    parts: List[str] = []
    for ev in events:
        if ev.get("type") == "user_input":
            parts.append(ev.get("text", ""))
        elif ev.get("type") == "assistant":
            content = (ev.get("message") or {}).get("content") or []
            for block in content:
                if block.get("type") == "text":
                    parts.append(block.get("text", ""))
    return "\n".join(parts)


class ChatRequest(BaseModel):
    message: str
    session_id: Optional[str] = None
    cwd: Optional[str] = None
    images: Optional[List[str]] = None
    model: Optional[str] = None
    system_prompt: Optional[str] = None
    display_message: Optional[str] = None
    permission_mode: Optional[str] = None
    allowed_tools: Optional[List[str]] = None
    disallowed_tools: Optional[List[str]] = None
    force_new: Optional[bool] = None
    # UI-only metadata for attached docs (name/size/length/path); rendered as
    # badges on the user message. Not used to build the prompt — the doc text
    # is already embedded in `message` by the client.
    docs: Optional[List[dict]] = None


class PromptRequest(BaseModel):
    name: str
    content: str


class SessionPatch(BaseModel):
    title: Optional[str] = None
    pinned: Optional[bool] = None
    archived: Optional[bool] = None
    tags: Optional[str] = None


class ForkRequest(BaseModel):
    event_index: int
    new_text: Optional[str] = None


class RestoreRequest(BaseModel):
    event_index: int


class FetchUrlRequest(BaseModel):
    url: str
    max_chars: Optional[int] = 10000


def build_args(
    message: str,
    session_id: str,
    resume: bool,
    model: Optional[str],
    system_prompt: Optional[str],
    permission_mode: Optional[str] = None,
    allowed_tools: Optional[List[str]] = None,
    disallowed_tools: Optional[List[str]] = None,
    use_stdin: bool = False,
) -> List[str]:
    args = ["claude"]
    if use_stdin:
        args += ["-p", "--input-format", "stream-json"]
    else:
        args += ["-p", message]
    args += [
        "--output-format", "stream-json",
        "--verbose",
        "--include-partial-messages",
    ]
    if resume:
        args += ["--resume", session_id]
    else:
        args += ["--session-id", session_id]
    if model:
        args += ["--model", model]
    if system_prompt:
        args += ["--append-system-prompt", system_prompt]
    if permission_mode and permission_mode in ("default", "acceptEdits", "bypassPermissions", "plan"):
        args += ["--permission-mode", permission_mode]
    if allowed_tools:
        args += ["--allowed-tools", ",".join(allowed_tools)]
    if disallowed_tools:
        args += ["--disallowed-tools", ",".join(disallowed_tools)]
    return args


def extract_tool_name(text: str) -> Optional[str]:
    mcp_match = re.search(r"\bmcp__[A-Za-z0-9_:-]+(?:__[A-Za-z0-9_:-]+)*\b", text)
    if mcp_match:
        return mcp_match.group(0)

    for tool in KNOWN_TOOL_NAMES:
        if re.search(rf"\b{re.escape(tool)}\b", text):
            return tool

    patterns = [
        r"(?:MCP tool|mcp tool|tool)\s+[\"'`]?([A-Za-z][A-Za-z0-9_:-]{1,80})[\"'`]?",
        r"[\"'`]([A-Za-z][A-Za-z0-9_:-]{1,80})[\"'`]\s+(?:tool|Tool|MCP tool|mcp tool)",
    ]
    stop_words = {"approval", "permission", "tool", "tools", "mcp", "required", "requires", "non-interactive"}
    for pattern in patterns:
        m = re.search(pattern, text)
        if not m:
            continue
        candidate = m.group(1).strip()
        if candidate.lower() not in stop_words:
            return candidate
    return None


def classify_claude_error(message: str) -> dict:
    text = (message or "").strip() or "claude exited with error"
    lower = text.lower()
    tool_name = extract_tool_name(text)
    permissionish = any(k in lower for k in (
        "requires approval", "approval required", "needs approval", "approval",
        "cannot prompt", "non-interactive", "not allowed", "permission denied",
    ))
    if ("permission" in lower and ("tool" in lower or "mcp" in lower or tool_name)) or (permissionish and ("tool" in lower or "mcp" in lower or tool_name)):
        return {
            "type": "permission_error",
            "message": text,
            "tool_name": tool_name,
            "hint": "当前 Web UI 不支持运行中批准工具权限；请预先放行工具后重试本轮，或改用 Claude Code CLI。",
        }
    return {"type": "error", "message": text}


def build_image_input_message(message: str, images: List[str]) -> bytes:
    """Build a stream-json user message. Works with or without images."""
    import base64 as b64mod
    content: List[dict] = []
    for img_path in images:
        p = Path(img_path)
        if not p.exists():
            continue
        ext = p.suffix.lower()
        media_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                     ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp"}
        media_type = media_map.get(ext, "image/png")
        data = b64mod.b64encode(p.read_bytes()).decode()
        content.append({"type": "image", "source": {"type": "base64", "media_type": media_type, "data": data}})
    content.append({"type": "text", "text": message})
    msg = {"type": "user", "message": {"role": "user", "content": content}}
    return json.dumps(msg, ensure_ascii=False).encode() + b"\n"


async def _git_run(cwd: str, *args: str) -> Optional[str]:
    try:
        proc = await asyncio.create_subprocess_exec(
            "git", "-C", cwd, *args,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=10)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            return None
        if proc.returncode != 0:
            return None
        return stdout.decode("utf-8", errors="replace").strip()
    except Exception:
        return None


async def create_git_checkpoint(cwd: str) -> Optional[dict]:
    if not cwd or not os.path.isdir(cwd):
        return None
    git_dir = await _git_run(cwd, "rev-parse", "--git-dir")
    if git_dir is None:
        return None
    head = await _git_run(cwd, "rev-parse", "HEAD")
    if head is None:
        return None
    stash = await _git_run(cwd, "stash", "create", f"claude-web-checkpoint-{int(time.time())}")
    return {"type": "git", "head": head, "stash": stash or ""}


async def restore_git_checkpoint(cwd: str, cp: dict) -> bool:
    if not cp or cp.get("type") != "git" or not cwd:
        return False
    head = cp.get("head")
    stash = cp.get("stash") or ""
    if not head:
        return False
    if await _git_run(cwd, "reset", "--hard", head) is None:
        return False
    await _git_run(cwd, "clean", "-fd")
    if stash:
        await _git_run(cwd, "stash", "apply", stash)
    return True


def format_context_snippet(events: List[dict], max_chars: int = 6000) -> str:
    lines: List[str] = []
    total = 0
    for ev in events:
        t = ev.get("type")
        if t == "user_input":
            text = (ev.get("text") or "").strip()
            if text:
                chunk = f"用户: {text}"
                lines.append(chunk)
                total += len(chunk)
        elif t == "assistant":
            content = (ev.get("message") or {}).get("content") or []
            for block in content:
                if block.get("type") == "text":
                    text = (block.get("text") or "").strip()
                    if text:
                        chunk = f"助手: {text[:600]}"
                        lines.append(chunk)
                        total += len(chunk)
                elif block.get("type") == "tool_use":
                    name = block.get("name", "")
                    chunk = f"(助手调用了工具: {name})"
                    lines.append(chunk)
                    total += len(chunk)
        if total > max_chars:
            lines.append("...（历史已截断）")
            break
    return "\n\n".join(lines)


def derive_title(message: str) -> str:
    """First line of the user message, with code fences / markdown headers / bullet
    markers stripped so the title reads naturally even when the message starts with
    a code block or markdown."""
    if not message:
        return "未命名会话"
    lines = message.splitlines()
    in_fence = False
    first_in_fence: Optional[str] = None
    for raw in lines:
        stripped = raw.strip()
        if not stripped:
            continue
        # Toggle on triple-backtick fences and skip their content.
        if stripped.startswith("```") or stripped.startswith("~~~"):
            in_fence = not in_fence
            continue
        if in_fence:
            if first_in_fence is None:
                first_in_fence = stripped
            continue
        # Strip markdown header / quote / list prefixes for nicer titles.
        cleaned = re.sub(r"^[#>\-*+\d.\s]+", "", stripped).strip()
        if cleaned:
            return cleaned[:60]
    if first_in_fence:
        return first_in_fence[:60]
    fallback = message.strip().replace("\n", " ")
    return fallback[:60] if fallback else "未命名会话"


def session_has_remote_conversation(events: List[dict]) -> bool:
    for ev in events:
        event_type = ev.get("type")
        if event_type == "assistant":
            return True
        if event_type == "system" and ev.get("subtype") == "init":
            return True
        if event_type == "result" and not ev.get("is_error"):
            return True
    return False


def resolve_remote_session_state(session_id: str, row: Optional[sqlite3.Row], events: List[dict]):
    if row is None:
        return session_id, session_has_remote_conversation(events)
    remote_session_id = (row["remote_session_id"] or "").strip() or session_id
    if (row["remote_session_id"] or "").strip():
        return remote_session_id, bool(row["remote_ready"])
    return remote_session_id, session_has_remote_conversation(events)


def set_session_remote_state(session_id: str, remote_session_id: str, remote_ready: bool) -> None:
    now = time.time()
    with db_connect() as conn:
        conn.execute(
            "UPDATE sessions SET remote_session_id = ?, remote_ready = ?, updated_at = ? WHERE id = ?",
            (remote_session_id, 1 if remote_ready else 0, now, session_id),
        )


@app.post("/api/chat")
async def chat(req: ChatRequest):
    session_id = req.session_id or str(uuid.uuid4())
    existing_events = load_events(session_id) if req.session_id else []
    with db_connect() as conn:
        row = conn.execute(
            "SELECT cwd, remote_session_id, remote_ready FROM sessions WHERE id = ?",
            (session_id,),
        ).fetchone()

    remote_session_id, remote_ready = resolve_remote_session_state(session_id, row, existing_events)
    if req.force_new is True:
        stored_remote_id = ((row["remote_session_id"] or "").strip() if row else "")
        stored_remote_ready = bool(row["remote_ready"]) if row else False
        if stored_remote_id and not stored_remote_ready:
            remote_session_id = stored_remote_id
        elif row is not None and remote_ready:
            remote_session_id = str(uuid.uuid4())
        remote_ready = False

    is_new = not remote_ready
    work_dir = req.cwd or (row["cwd"] if row and row["cwd"] else os.path.expanduser("~"))
    full_message = req.message
    display_text = req.display_message if req.display_message is not None else req.message

    checkpoint = await create_git_checkpoint(work_dir)

    user_event = {
        "type": "user_input",
        "text": display_text,
        "images": req.images or [],
        "docs": req.docs or [],
        "ts": time.time(),
        "checkpoint": checkpoint,
    }
    # When the prompt was rewritten on the client (doc content / URL fetch / web-search prefix
    # injected), keep the full sent text so badge previews can recover doc bodies even
    # after the upload file is pruned. Only stored when it actually differs.
    if req.message != display_text:
        user_event["full_text"] = req.message
    append_event(session_id, user_event)
    upsert_session(session_id, derive_title(display_text), work_dir)
    set_session_remote_state(session_id, remote_session_id, remote_ready and not is_new)

    async def generate():
        remote_became_ready = remote_ready and not is_new
        meta = {
            "type": "meta",
            "session_id": session_id,
            "cwd": work_dir,
            "has_checkpoint": checkpoint is not None,
        }
        yield f"data: {json.dumps(meta)}\n\n"

        # If a previous chat for the same session is still running (e.g. duplicate
        # request, network retry, fast double-click), terminate it before spawning
        # a new one. Otherwise the old subprocess would be orphaned, burning tokens
        # and producing stray events.
        existing = _running_processes.pop(session_id, None)
        if existing is not None:
            await _terminate_process(existing)

        has_images = bool(req.images)
        # Route through stdin when the prompt would blow past argv limits
        # (macOS ~256KB, Linux ~128KB total argv). Images already force stdin.
        message_too_large = len(full_message.encode("utf-8")) > ARGV_STDIN_THRESHOLD
        use_stdin = has_images or message_too_large
        args = build_args(
            full_message, remote_session_id,
            resume=not is_new,
            model=req.model,
            system_prompt=req.system_prompt,
            permission_mode=req.permission_mode,
            allowed_tools=req.allowed_tools,
            disallowed_tools=req.disallowed_tools,
            use_stdin=use_stdin,
        )
        stdin_data: Optional[bytes] = None
        if use_stdin:
            stdin_data = build_image_input_message(full_message, req.images or [])
        try:
            process = await asyncio.create_subprocess_exec(
                *args,
                stdin=asyncio.subprocess.PIPE if use_stdin else None,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                cwd=work_dir,
                limit=16 * 1024 * 1024,
            )
            if use_stdin and stdin_data and process.stdin:
                try:
                    process.stdin.write(stdin_data)
                    await process.stdin.drain()
                    process.stdin.close()
                    await process.stdin.wait_closed()
                except (BrokenPipeError, ConnectionResetError):
                    # CLI exited early (auth failure, bad args, etc). The stderr
                    # path will surface the real reason; don't tear down SSE here.
                    pass
        except FileNotFoundError:
            err_event = {"type": "error", "message": "claude CLI not found in PATH"}
            append_event(session_id, err_event)
            yield f"data: {json.dumps(err_event)}\n\n"
            return

        _running_processes[session_id] = process
        stderr_buffer = bytearray()
        stderr_task: Optional[asyncio.Task] = None
        if process.stderr is not None:
            stderr_task = asyncio.create_task(_drain_stream(process.stderr, stderr_buffer))

        try:
            assert process.stdout is not None
            while True:
                try:
                    raw = await process.stdout.readline()
                except ValueError as e:
                    err_event = {"type": "error", "message": f"stdout line too large: {e}"}
                    append_event(session_id, err_event)
                    yield f"data: {json.dumps(err_event)}\n\n"
                    break
                if not raw:
                    break
                line = raw.decode("utf-8", errors="replace").strip()
                if not line:
                    continue
                try:
                    obj = json.loads(line)
                except json.JSONDecodeError:
                    obj = {"type": "raw", "text": line}
                t = obj.get("type")
                if session_has_remote_conversation([obj]):
                    remote_became_ready = True
                if t != "stream_event" and not (t == "system" and obj.get("subtype", "").startswith("hook_")):
                    append_event(session_id, obj)
                yield f"data: {json.dumps(obj, ensure_ascii=False)}\n\n"

            rc = await process.wait()
            if stderr_task is not None:
                try:
                    await asyncio.wait_for(asyncio.shield(stderr_task), timeout=1.0)
                except asyncio.TimeoutError:
                    pass
            stopped_by_user = session_id in _stopped_sessions
            if rc != 0 and not stopped_by_user:
                err_text = bytes(stderr_buffer).decode("utf-8", errors="replace")
                err_event = classify_claude_error(err_text or f"claude exited with code {rc}")
                append_event(session_id, err_event)
                yield f"data: {json.dumps(err_event, ensure_ascii=False)}\n\n"
        finally:
            if stderr_task is not None and not stderr_task.done():
                stderr_task.cancel()
                try:
                    await stderr_task
                except (asyncio.CancelledError, Exception):
                    pass
            await _terminate_process(process)
            _running_processes.pop(session_id, None)
            _stopped_sessions.discard(session_id)

        upsert_session(session_id, derive_title(display_text), work_dir)
        if remote_became_ready:
            set_session_remote_state(session_id, remote_session_id, True)
        yield f"data: {json.dumps({'type': 'done'})}\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.post("/api/chat/stop/{session_id}")
async def stop_chat(session_id: str):
    process = _running_processes.get(session_id)
    if process is None:
        raise HTTPException(status_code=404, detail="no running process for this session")
    _stopped_sessions.add(session_id)
    await _terminate_process(process)
    stop_event = {"type": "error", "message": "用户中止", "ts": time.time()}
    append_event(session_id, stop_event)
    return {"ok": True}


@app.post("/api/sessions/{session_id}/prepare-fork")
async def prepare_fork(session_id: str, req: ForkRequest):
    events = load_events(session_id)
    user_event_positions = [i for i, e in enumerate(events) if e.get("type") == "user_input"]
    if req.event_index < 0 or req.event_index >= len(user_event_positions):
        raise HTTPException(status_code=400, detail="invalid event_index")

    target_pos = user_event_positions[req.event_index]
    events_before = events[:target_pos]
    original_text = events[target_pos].get("text", "")
    new_text = req.new_text if req.new_text is not None and req.new_text.strip() else original_text

    with db_connect() as conn:
        row = conn.execute("SELECT cwd FROM sessions WHERE id = ?", (session_id,)).fetchone()
    cwd = row["cwd"] if row else os.path.expanduser("~")

    new_id = str(uuid.uuid4())
    upsert_session(new_id, derive_title(new_text), cwd)

    with db_connect() as conn:
        conn.execute(
            "UPDATE sessions SET tags = ? WHERE id = ?",
            (f"forked-from-{session_id[:8]}", new_id),
        )

    context = format_context_snippet(events_before)
    if context:
        packed_message = (
            "【以下是之前的对话历史，仅作为参考上下文（不要重复回应历史问题）】\n"
            f"{context}\n\n"
            "【请基于以上历史上下文，回应这个新问题】\n"
            f"{new_text}"
        )
    else:
        packed_message = new_text

    return {
        "session_id": new_id,
        "cwd": cwd,
        "sent_message": packed_message,
        "display_message": new_text,
        "forked_from": session_id,
    }


@app.post("/api/sessions/{session_id}/prepare-inline-edit")
async def prepare_inline_edit(session_id: str, req: ForkRequest):
    if session_id in _running_processes:
        raise HTTPException(status_code=409, detail="session is running")

    events = load_events(session_id)
    user_event_positions = [i for i, e in enumerate(events) if e.get("type") == "user_input"]
    if req.event_index < 0 or req.event_index >= len(user_event_positions):
        raise HTTPException(status_code=400, detail="invalid event_index")

    target_pos = user_event_positions[req.event_index]
    events_before = events[:target_pos]
    original_event = events[target_pos]
    original_text = original_event.get("text", "")
    original_images = original_event.get("images", []) or []
    new_text = req.new_text if req.new_text is not None and req.new_text.strip() else original_text

    with db_connect() as conn:
        row = conn.execute("SELECT cwd FROM sessions WHERE id = ?", (session_id,)).fetchone()
    cwd = row["cwd"] if row else os.path.expanduser("~")

    save_events(session_id, events_before)
    upsert_session(session_id, derive_title(new_text), cwd)
    set_session_remote_state(session_id, str(uuid.uuid4()), False)

    context = format_context_snippet(events_before)
    if context:
        packed_message = (
            "【以下是之前的对话历史，仅作为参考上下文（不要重复回应历史问题）】\n"
            f"{context}\n\n"
            "【请基于以上历史上下文，继续这个对话，并回应下面这条经过编辑的新消息】\n"
            f"{new_text}"
        )
    else:
        packed_message = new_text

    return {
        "session_id": session_id,
        "cwd": cwd,
        "sent_message": packed_message,
        "display_message": new_text,
        "images": original_images,
    }


@app.post("/api/sessions/{session_id}/restore-checkpoint")
async def restore_checkpoint(session_id: str, req: RestoreRequest):
    events = load_events(session_id)
    user_event_positions = [i for i, e in enumerate(events) if e.get("type") == "user_input"]
    if req.event_index < 0 or req.event_index >= len(user_event_positions):
        raise HTTPException(status_code=400, detail="invalid event_index")
    ev = events[user_event_positions[req.event_index]]
    cp = ev.get("checkpoint")
    if not cp:
        raise HTTPException(status_code=400, detail="no checkpoint on this turn")

    with db_connect() as conn:
        row = conn.execute("SELECT cwd FROM sessions WHERE id = ?", (session_id,)).fetchone()
    cwd = row["cwd"] if row else ""
    if not cwd:
        raise HTTPException(status_code=400, detail="session has no cwd")

    ok = await restore_git_checkpoint(cwd, cp)
    if not ok:
        raise HTTPException(status_code=500, detail="restore failed")
    return {"ok": True, "cwd": cwd}


@app.post("/api/upload")
async def upload(file: UploadFile = File(...)):
    if file.filename is None:
        raise HTTPException(status_code=400, detail="filename missing")
    ext = Path(file.filename).suffix.lower()
    if ext not in IMAGE_EXTS:
        raise HTTPException(status_code=400, detail=f"unsupported type {ext}")

    data = await file.read()
    if len(data) > MAX_UPLOAD_MB * 1024 * 1024:
        raise HTTPException(status_code=413, detail=f"file exceeds {MAX_UPLOAD_MB} MB")

    name = f"{uuid.uuid4().hex}{ext}"
    path = UPLOADS_DIR / name
    path.write_bytes(data)

    return {
        "path": str(path.absolute()),
        "url": f"/uploads/{name}",
        "name": file.filename,
        "size": len(data),
    }


DOC_MIME_EXTS = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.ms-excel.sheet.macroenabled.12": ".xlsm",
    "application/vnd.ms-excel": ".xls",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/xhtml+xml": ".html",
    "application/javascript": ".js",
    "application/json": ".json",
    "application/xml": ".xml",
    "image/svg+xml": ".svg",
    "text/csv": ".csv",
    "text/css": ".css",
    "text/html": ".html",
    "text/javascript": ".js",
    "text/tab-separated-values": ".tsv",
    "text/markdown": ".md",
    "text/plain": ".txt",
    "text/xml": ".xml",
}
MAX_DOC_MB = 30
# Soft cap kept for UI display; we no longer hard-truncate the document text on
# upload. Anything beyond this just gets a "large document" hint in the response.
LARGE_DOC_CHARS_HINT = 200_000
# Argv length safety margin. macOS allows ~256KB total argv; once the prompt
# (UTF-8 bytes) crosses this we route through stdin to avoid E2BIG.
ARGV_STDIN_THRESHOLD = 60_000


def _extract_pdf_text(path: Path) -> str:
    """Extract PDF text. Prefers pdfplumber (better tables/layout) when available,
    falls back to pypdf on any failure (import miss, malformed PDF, table extraction error).
    Each page is prefixed with [Page N] so the model can cite."""
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        pdfplumber = None  # type: ignore
    if pdfplumber is not None:
        try:
            parts: List[str] = []
            with pdfplumber.open(str(path)) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text() or ""
                    tables = []
                    try:
                        for table in page.extract_tables() or []:
                            if not table:
                                continue
                            rows = [" | ".join((cell or "").strip() for cell in row) for row in table]
                            tables.append("\n".join(rows))
                    except Exception:
                        pass
                    section = f"[Page {i}]\n{page_text}"
                    if tables:
                        section += "\n\n" + "\n\n".join(tables)
                    parts.append(section)
            return "\n\n".join(parts)
        except Exception:
            # Any pdfplumber failure (malformed PDF, missing deps, parse error) → fall through to pypdf.
            pass
    import pypdf
    reader = pypdf.PdfReader(str(path))
    parts = []
    for i, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        parts.append(f"[Page {i}]\n{text}")
    return "\n\n".join(parts)


def _docx_table_to_markdown(table) -> str:
    rows = []
    for row in table.rows:
        cells = [(cell.text or "").strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if not rows:
        return ""
    if len(rows) == 1:
        return rows[0]
    header_sep = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
    return rows[0] + "\n" + header_sep + "\n" + "\n".join(rows[1:])


def _docx_hf_lines(hdr_or_ftr, label: str) -> List[str]:
    """Collect paragraphs and tables from a header/footer container as labeled lines."""
    if hdr_or_ftr is None:
        return []
    out: List[str] = []
    for p in hdr_or_ftr.paragraphs:
        if p.text and p.text.strip():
            out.append(f"[{label}] {p.text}")
    for t in getattr(hdr_or_ftr, "tables", []) or []:
        md = _docx_table_to_markdown(t)
        if md:
            out.append(f"[{label} table]\n{md}")
    return out


def _extract_docx_text(path: Path) -> str:
    """Extract DOCX content preserving the original paragraph/table order.
    Walks the body XML in document order so a 'paragraph → table → paragraph' layout
    survives instead of becoming 'all paragraphs then all tables'.
    Includes default / first-page / even-page headers and footers, plus any
    tables embedded inside them."""
    import docx
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph
    from docx.table import Table
    doc = docx.Document(str(path))
    parts: List[str] = []
    for section in doc.sections:
        parts += _docx_hf_lines(section.header, "Header")
        parts += _docx_hf_lines(section.first_page_header, "Header (first page)")
        parts += _docx_hf_lines(getattr(section, "even_page_header", None), "Header (even page)")
    body = doc.element.body
    para_tag = qn("w:p")
    table_tag = qn("w:tbl")
    for child in body.iterchildren():
        if child.tag == para_tag:
            p = Paragraph(child, doc)
            if p.text:
                parts.append(p.text)
        elif child.tag == table_tag:
            t = Table(child, doc)
            md = _docx_table_to_markdown(t)
            if md:
                parts.append(md)
    for section in doc.sections:
        parts += _docx_hf_lines(section.footer, "Footer")
        parts += _docx_hf_lines(section.first_page_footer, "Footer (first page)")
        parts += _docx_hf_lines(getattr(section, "even_page_footer", None), "Footer (even page)")
    return "\n".join(parts)


def _extract_xlsx_text(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    try:
        parts: List[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    parts.append(" | ".join(cells))
    finally:
        wb.close()
    return "\n".join(parts)


def _extract_xls_text(path: Path) -> str:
    import xlrd
    wb = xlrd.open_workbook(str(path), on_demand=True)
    try:
        parts: List[str] = []
        for sheet in wb.sheets():
            parts.append(f"[Sheet: {sheet.name}]")
            for row_idx in range(sheet.nrows):
                values = []
                for cell in sheet.row(row_idx):
                    value = cell.value
                    if isinstance(value, float) and value.is_integer():
                        value = int(value)
                    values.append(str(value) if value != "" else "")
                if any(values):
                    parts.append(" | ".join(values))
    finally:
        wb.release_resources()
    return "\n".join(parts)


def _extract_pptx_text(path: Path) -> str:
    """Extract PowerPoint slides as plain text. Each slide gets a [Slide N] header
    so the model can cite. Pulls title, body text from every shape (recursing into
    grouped shapes), embedded tables (markdown-ified), and the speaker notes pane."""
    from pptx import Presentation  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore

    def walk(shape, title_shape, body_lines: List[str]) -> None:
        # Recurse into groups so text/tables nested inside a group aren't lost.
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                walk(child, title_shape, body_lines)
            return
        if title_shape is not None and shape == title_shape:
            return
        if shape.has_text_frame:
            text = (shape.text_frame.text or "").strip()
            if text:
                body_lines.append(text)
        elif shape.has_table:
            rows = []
            for row in shape.table.rows:
                cells = [(c.text or "").strip().replace("\n", " ") for c in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                if len(rows) > 1:
                    sep = "| " + " | ".join("---" for _ in shape.table.rows[0].cells) + " |"
                    body_lines.append(rows[0] + "\n" + sep + "\n" + "\n".join(rows[1:]))
                else:
                    body_lines.append(rows[0])

    prs = Presentation(str(path))
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        title_shape = None
        title = ""
        try:
            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.has_text_frame:
                title = (title_shape.text_frame.text or "").strip()
        except Exception:
            title_shape = None
        header = f"[Slide {i}]" + (f" {title}" if title else "")
        body_lines: List[str] = []
        for shape in slide.shapes:
            try:
                walk(shape, title_shape, body_lines)
            except Exception:
                continue  # don't let one bad shape kill the whole slide
        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            pass
        section = header
        if body_lines:
            section += "\n" + "\n".join(body_lines)
        if notes_text:
            section += f"\n[Notes] {notes_text}"
        parts.append(section)
    return "\n\n".join(parts)


def _looks_binary(data: bytes) -> bool:
    sample = data[:8192]
    if not sample:
        return False
    if b"\x00" in sample:
        return True
    allowed_controls = {9, 10, 12, 13}
    control_count = sum(1 for b in sample if b < 32 and b not in allowed_controls)
    return control_count / len(sample) > 0.30


def _reject_mojibake(text: str) -> str:
    if not text:
        return text
    replacement_count = text.count("\ufffd")
    if replacement_count and replacement_count / len(text) > 0.01:
        raise HTTPException(status_code=400, detail="unsupported binary file")
    return text


def _decode_text_upload(data: bytes) -> str:
    if data.startswith((b"\xff\xfe", b"\xfe\xff")):
        return _reject_mojibake(data.decode("utf-16", errors="replace"))
    if data.startswith(b"\xef\xbb\xbf"):
        return _reject_mojibake(data.decode("utf-8-sig", errors="replace"))
    if _looks_binary(data):
        raise HTTPException(status_code=400, detail="unsupported binary file")

    for encoding in ("utf-8", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return _reject_mojibake(data.decode("utf-8", errors="replace"))


def _doc_ext_from_upload(file: UploadFile) -> tuple[str, str]:
    filename = file.filename or "clipboard-file"
    ext = Path(filename).suffix.lower()
    if not ext:
        content_type = (file.content_type or "").split(";", 1)[0].strip().lower()
        ext = DOC_MIME_EXTS.get(content_type, "")
    return ext, filename


@app.post("/api/upload-doc")
async def upload_doc(file: UploadFile = File(...)):
    ext, filename = _doc_ext_from_upload(file)
    data = await file.read()
    if len(data) > MAX_DOC_MB * 1024 * 1024:
        raise HTTPException(status_code=413, detail=f"file exceeds {MAX_DOC_MB} MB")

    name = f"{uuid.uuid4().hex}{ext}"
    path = UPLOADS_DIR / name
    path.write_bytes(data)

    try:
        if ext == ".pdf":
            text = _extract_pdf_text(path)
        elif ext == ".docx":
            text = _extract_docx_text(path)
        elif ext == ".pptx":
            text = _extract_pptx_text(path)
        elif ext in (".xlsx", ".xlsm"):
            text = _extract_xlsx_text(path)
        elif ext == ".xls":
            text = _extract_xls_text(path)
        elif ext in (".html", ".htm", ".xhtml"):
            text = _extract_html_text(_decode_text_upload(data))
        else:
            text = _decode_text_upload(data)
    except HTTPException:
        path.unlink(missing_ok=True)
        raise
    except Exception as e:
        path.unlink(missing_ok=True)
        raise HTTPException(status_code=500, detail=f"extract failed: {e}")

    text = text.strip()
    is_large = len(text) > LARGE_DOC_CHARS_HINT

    return {
        "path": str(path.absolute()),
        "name": filename,
        "size": len(data),
        "ext": ext,
        "content": text,
        "length": len(text),
        "truncated": False,
        "large": is_large,
    }


@app.get("/api/doc-content")
async def doc_content(path: str = Query(...)):
    """Read back the extracted text for a doc badge preview.
    Locked to files inside UPLOADS_DIR to prevent path traversal."""
    try:
        target = Path(path).resolve()
    except (OSError, ValueError):
        raise HTTPException(status_code=400, detail="invalid path")
    uploads_root = UPLOADS_DIR.resolve()
    try:
        target.relative_to(uploads_root)
    except ValueError:
        raise HTTPException(status_code=403, detail="path outside uploads directory")
    if not target.is_file():
        raise HTTPException(status_code=404, detail="file not found")

    ext = target.suffix.lower()
    try:
        if ext == ".pdf":
            text = _extract_pdf_text(target)
        elif ext == ".docx":
            text = _extract_docx_text(target)
        elif ext == ".pptx":
            text = _extract_pptx_text(target)
        elif ext in (".xlsx", ".xlsm"):
            text = _extract_xlsx_text(target)
        elif ext == ".xls":
            text = _extract_xls_text(target)
        elif ext in (".html", ".htm", ".xhtml"):
            text = _extract_html_text(_decode_text_upload(target.read_bytes()))
        else:
            text = _decode_text_upload(target.read_bytes())
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"read failed: {e}")
    return {"content": text.strip(), "length": len(text)}


class ExecCodeRequest(BaseModel):
    language: str
    code: str
    timeout: Optional[int] = 10


EXEC_LANG_MAP: Dict[str, List[str]] = {
    "python": ["python3", "-c"],
    "python3": ["python3", "-c"],
    "py": ["python3", "-c"],
    "javascript": ["node", "-e"],
    "js": ["node", "-e"],
    "node": ["node", "-e"],
    "bash": ["bash", "-c"],
    "sh": ["bash", "-c"],
    "shell": ["bash", "-c"],
}


@app.post("/api/exec-code")
async def exec_code(req: ExecCodeRequest):
    lang = (req.language or "").lower().strip()
    cmd = EXEC_LANG_MAP.get(lang)
    if cmd is None:
        raise HTTPException(status_code=400, detail=f"unsupported language: {lang}")
    if not req.code or len(req.code) > 100_000:
        raise HTTPException(status_code=400, detail="code empty or too large")

    timeout = max(1, min(int(req.timeout or 10), 30))
    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd, req.code,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=str(UPLOADS_DIR),
        )
        try:
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
        except asyncio.TimeoutError:
            proc.kill()
            return {"stdout": "", "stderr": f"execution timed out after {timeout}s", "returncode": -1, "timed_out": True}
    except FileNotFoundError as e:
        raise HTTPException(status_code=500, detail=f"interpreter not found: {e}")

    return {
        "stdout": stdout.decode("utf-8", errors="replace")[:50_000],
        "stderr": stderr.decode("utf-8", errors="replace")[:10_000],
        "returncode": proc.returncode,
        "timed_out": False,
    }


def _row_to_session(r: sqlite3.Row) -> dict:
    tags = [t for t in (r["tags"] or "").split(",") if t]
    return {
        "id": r["id"],
        "title": r["title"] or "未命名会话",
        "cwd": r["cwd"],
        "created_at": r["created_at"],
        "updated_at": r["updated_at"],
        "pinned": bool(r["pinned"]),
        "archived": bool(r["archived"]),
        "tags": tags,
    }


@app.get("/api/sessions")
async def list_sessions(q: Optional[str] = None, archived: bool = False, tag: Optional[str] = None):
    with db_connect() as conn:
        where = "archived = 1" if archived else "archived = 0"
        rows = conn.execute(
            f"SELECT id, title, cwd, created_at, updated_at, pinned, archived, tags FROM sessions "
            f"WHERE {where} ORDER BY pinned DESC, updated_at DESC LIMIT 500"
        ).fetchall()

    items = [_row_to_session(r) for r in rows]

    if tag:
        items = [i for i in items if tag in i["tags"]]

    if q:
        q_lower = q.lower()
        filtered: List[dict] = []
        for item in items:
            if q_lower in item["title"].lower() or q_lower in ",".join(item["tags"]).lower():
                filtered.append(item)
                continue
            try:
                events = load_events(item["id"])
                content = summarize_text_from_events(events).lower()
                if q_lower in content:
                    filtered.append(item)
            except Exception:
                continue
        items = filtered

    return items


@app.get("/api/sessions/{session_id}")
async def get_session(session_id: str):
    with db_connect() as conn:
        row = conn.execute(
            "SELECT id, title, cwd, created_at, updated_at, pinned, archived, tags FROM sessions WHERE id = ?",
            (session_id,),
        ).fetchone()
    if row is None:
        raise HTTPException(status_code=404, detail="session not found")
    data = _row_to_session(row)
    data["events"] = load_events(session_id)
    return data


@app.patch("/api/sessions/{session_id}")
async def patch_session(session_id: str, req: SessionPatch):
    updates: List[str] = []
    params: List = []
    if req.title is not None:
        updates += ["title = ?", "manual_title = 1"]
        params.append(req.title)
    if req.pinned is not None:
        updates.append("pinned = ?")
        params.append(1 if req.pinned else 0)
    if req.archived is not None:
        updates.append("archived = ?")
        params.append(1 if req.archived else 0)
    if req.tags is not None:
        updates.append("tags = ?")
        params.append(req.tags)
    if not updates:
        return {"ok": True}
    params.append(session_id)
    with db_connect() as conn:
        conn.execute(f"UPDATE sessions SET {', '.join(updates)} WHERE id = ?", params)
    return {"ok": True}


@app.delete("/api/sessions/{session_id}")
async def delete_session(session_id: str):
    with db_connect() as conn:
        conn.execute("DELETE FROM sessions WHERE id = ?", (session_id,))
    path = HISTORY_DIR / f"{session_id}.jsonl"
    if path.exists():
        path.unlink()
    return {"ok": True}


@app.post("/api/sessions/{session_id}/suggest-title")
async def suggest_title(session_id: str):
    events = load_events(session_id)
    if not events:
        raise HTTPException(status_code=404, detail="empty session")
    summary = summarize_text_from_events(events)[:3000]
    if not summary.strip():
        raise HTTPException(status_code=400, detail="no textual content")
    prompt = f"根据下面的对话，用中文生成一个不超过15字、不带引号的会话标题（只输出标题本身）：\n\n{summary}"
    try:
        proc = await asyncio.create_subprocess_exec(
            "claude", "-p", prompt, "--output-format", "text",
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=60)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            raise HTTPException(status_code=504, detail="title generation timeout")
    except HTTPException:
        raise
    title = stdout.decode("utf-8", errors="replace").strip().splitlines()[0].strip(' "\'"""''').strip()[:60]
    if not title:
        raise HTTPException(status_code=500, detail="empty title")
    with db_connect() as conn:
        conn.execute("UPDATE sessions SET title = ?, manual_title = 1 WHERE id = ?", (title, session_id))
    return {"title": title}


@app.get("/api/sessions/{session_id}/export")
async def export_session(session_id: str):
    events = load_events(session_id)
    if not events:
        raise HTTPException(status_code=404, detail="session not found")
    with db_connect() as conn:
        row = conn.execute("SELECT title FROM sessions WHERE id = ?", (session_id,)).fetchone()
    title = row["title"] if row and row["title"] else session_id

    lines: List[str] = [f"# {title}", "", f"_会话 ID: {session_id}_", ""]
    for ev in events:
        t = ev.get("type")
        if t == "user_input":
            lines += ["## 👤 用户", "", ev.get("text", "")]
            for img in ev.get("images", []) or []:
                lines.append(f"![image]({img})")
            lines.append("")
        elif t == "assistant":
            content = (ev.get("message") or {}).get("content") or []
            for block in content:
                if block.get("type") == "text":
                    lines += ["## 🤖 Claude", "", block.get("text", ""), ""]
                elif block.get("type") == "tool_use":
                    name = block.get("name", "?")
                    lines += [f"### 🔧 工具调用: `{name}`", "", "```json",
                              json.dumps(block.get("input", {}), ensure_ascii=False, indent=2), "```", ""]
        elif t == "user":
            content = (ev.get("message") or {}).get("content") or []
            for block in content:
                if block.get("type") == "tool_result":
                    ct = block.get("content", "")
                    if isinstance(ct, list):
                        ct = "\n".join(c.get("text", "") if isinstance(c, dict) else str(c) for c in ct)
                    lines += ["### 📋 工具结果", "", "```", str(ct)[:5000], "```", ""]

    md = "\n".join(lines)
    return Response(
        content=md,
        media_type="text/markdown; charset=utf-8",
        headers={"Content-Disposition": f'attachment; filename="{session_id}.md"'},
    )


@app.get("/api/prompts")
async def list_prompts():
    with db_connect() as conn:
        rows = conn.execute(
            "SELECT id, name, content, created_at FROM prompts ORDER BY created_at DESC"
        ).fetchall()
    return [
        {"id": r["id"], "name": r["name"], "content": r["content"], "created_at": r["created_at"]}
        for r in rows
    ]


@app.post("/api/prompts")
async def create_prompt(req: PromptRequest):
    pid = uuid.uuid4().hex
    with db_connect() as conn:
        conn.execute(
            "INSERT INTO prompts (id, name, content, created_at) VALUES (?, ?, ?, ?)",
            (pid, req.name, req.content, time.time()),
        )
    return {"id": pid}


@app.put("/api/prompts/{prompt_id}")
async def update_prompt(prompt_id: str, req: PromptRequest):
    with db_connect() as conn:
        conn.execute(
            "UPDATE prompts SET name = ?, content = ? WHERE id = ?",
            (req.name, req.content, prompt_id),
        )
    return {"ok": True}


@app.delete("/api/prompts/{prompt_id}")
async def delete_prompt(prompt_id: str):
    with db_connect() as conn:
        conn.execute("DELETE FROM prompts WHERE id = ?", (prompt_id,))
    return {"ok": True}


@app.post("/api/suggest-followups")
async def suggest_followups(session_id: str = ""):
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id required")
    events = load_events(session_id)
    if not events:
        return {"suggestions": []}
    snippet = summarize_text_from_events(events[-20:])[-3000:]
    if not snippet.strip():
        return {"suggestions": []}
    prompt = (
        "根据以下对话内容，生成3个用户可能想继续追问的简短问题（每个不超过20字）。"
        "只输出3行，每行一个问题，不要编号、不要引号、不要其他内容。\n\n"
        f"{snippet}"
    )
    try:
        proc = await asyncio.create_subprocess_exec(
            "claude", "-p", prompt, "--output-format", "text", "--model", "haiku",
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=30)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            return {"suggestions": []}
    except Exception:
        return {"suggestions": []}
    lines = [l.strip() for l in stdout.decode("utf-8", errors="replace").splitlines() if l.strip()]
    suggestions = [l.lstrip("0123456789.-、）) ") for l in lines[:3]]
    return {"suggestions": suggestions}


# ===== MCP Management =====

_CLAUDE_CONFIG_PATH = Path.home() / ".claude.json"
_PROJECT_MCP_FILENAME = ".mcp.json"
_DISABLED_MCP_SERVERS_KEY = "claudeWebDisabledMcpServers"
_MCP_SCOPES = {"local", "user", "project"}


def _read_json_object(path: Path) -> dict:
    if not path.exists():
        return {}
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"invalid JSON in {path}: {e.msg}")
    except OSError as e:
        raise HTTPException(status_code=500, detail=f"cannot read {path}: {e}")
    if not isinstance(data, dict):
        raise HTTPException(status_code=500, detail=f"invalid JSON object in {path}")
    return data


def _write_json_object(path: Path, data: dict) -> None:
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
    except OSError as e:
        raise HTTPException(status_code=500, detail=f"cannot write {path}: {e}")


def _normalize_mcp_scope(scope: Optional[str]) -> str:
    normalized = (scope or "local").strip().lower()
    if normalized not in _MCP_SCOPES:
        raise HTTPException(status_code=400, detail="scope must be local, user, or project")
    return normalized


def _resolve_mcp_cwd(cwd: Optional[str]) -> Path:
    raw = (cwd or "").strip() or "~"
    target = Path(os.path.expanduser(raw)).resolve()
    if not target.is_dir():
        raise HTTPException(status_code=400, detail=f"invalid cwd: {raw}")
    return target


def _dict_value(parent: dict, key: str) -> dict:
    value = parent.get(key)
    if not isinstance(value, dict):
        value = {}
        parent[key] = value
    return value


def _source_label(scope: str) -> str:
    return {
        "local": "Local",
        "user": "User",
        "project": "Project",
    }.get(scope, scope)


def _mcp_sources(cwd: Optional[str], create: bool = False) -> list[dict]:
    project_dir = _resolve_mcp_cwd(cwd)
    project_key = str(project_dir)
    claude_data = _read_json_object(_CLAUDE_CONFIG_PATH)
    project_mcp_path = project_dir / _PROJECT_MCP_FILENAME
    project_mcp_data = _read_json_object(project_mcp_path)

    projects = _dict_value(claude_data, "projects") if create else claude_data.get("projects", {})
    if not isinstance(projects, dict):
        projects = {}
    local_project = projects.get(project_key)
    if create:
        local_project = projects.setdefault(project_key, {})
    if not isinstance(local_project, dict):
        local_project = {}

    project_choice = local_project if isinstance(local_project, dict) else {}
    disabled_project_servers = project_choice.get("disabledMcpjsonServers", [])
    if not isinstance(disabled_project_servers, list):
        disabled_project_servers = []

    return [
        {
            "scope": "user",
            "path": _CLAUDE_CONFIG_PATH,
            "data": claude_data,
            "servers": claude_data.get("mcpServers", {}) if isinstance(claude_data.get("mcpServers", {}), dict) else {},
            "disabled_servers": claude_data.get(_DISABLED_MCP_SERVERS_KEY, {}) if isinstance(claude_data.get(_DISABLED_MCP_SERVERS_KEY, {}), dict) else {},
        },
        {
            "scope": "local",
            "path": _CLAUDE_CONFIG_PATH,
            "data": claude_data,
            "servers": local_project.get("mcpServers", {}) if isinstance(local_project.get("mcpServers", {}), dict) else {},
            "disabled_servers": local_project.get(_DISABLED_MCP_SERVERS_KEY, {}) if isinstance(local_project.get(_DISABLED_MCP_SERVERS_KEY, {}), dict) else {},
            "project_key": project_key,
            "project": local_project,
        },
        {
            "scope": "project",
            "path": project_mcp_path,
            "data": project_mcp_data,
            "servers": project_mcp_data.get("mcpServers", {}) if isinstance(project_mcp_data.get("mcpServers", {}), dict) else {},
            "disabled_names": set(str(v) for v in disabled_project_servers),
            "claude_data": claude_data,
            "project_key": project_key,
            "project": local_project,
        },
    ]


def _mcp_target(scope: str, cwd: Optional[str]) -> dict:
    normalized = _normalize_mcp_scope(scope)
    sources = _mcp_sources(cwd, create=True)
    for source in sources:
        if source["scope"] == normalized:
            if normalized == "user":
                source["servers"] = _dict_value(source["data"], "mcpServers")
                source["disabled_servers"] = _dict_value(source["data"], _DISABLED_MCP_SERVERS_KEY)
            elif normalized == "local":
                source["servers"] = _dict_value(source["project"], "mcpServers")
                source["disabled_servers"] = _dict_value(source["project"], _DISABLED_MCP_SERVERS_KEY)
            else:
                source["servers"] = _dict_value(source["data"], "mcpServers")
                disabled = source["project"].get("disabledMcpjsonServers")
                if not isinstance(disabled, list):
                    disabled = []
                    source["project"]["disabledMcpjsonServers"] = disabled
                source["disabled_names"] = set(str(v) for v in disabled)
            return source
    raise HTTPException(status_code=400, detail="invalid scope")


def _save_mcp_source(source: dict, save_claude_choices: bool = False) -> None:
    _write_json_object(source["path"], source["data"])
    if source["scope"] == "project" and save_claude_choices:
        _write_json_object(_CLAUDE_CONFIG_PATH, source["claude_data"])


def _find_mcp_source(name: str, scope: Optional[str], cwd: Optional[str]) -> dict:
    if scope:
        source = _mcp_target(scope, cwd)
        if _mcp_config_in_source(source, name) is None:
            raise HTTPException(status_code=404, detail=f"server '{name}' not found")
        return source

    matches = []
    for source in _mcp_sources(cwd, create=True):
        if _mcp_config_in_source(source, name) is not None:
            matches.append(source)
    if not matches:
        raise HTTPException(status_code=404, detail=f"server '{name}' not found")
    if len(matches) > 1:
        scopes = ", ".join(_source_label(m["scope"]) for m in matches)
        raise HTTPException(status_code=409, detail=f"server '{name}' exists in multiple scopes: {scopes}")
    return matches[0]


def _mcp_config_in_source(source: dict, name: str) -> Optional[dict]:
    servers = source.get("servers") or {}
    if name in servers and isinstance(servers[name], dict):
        return servers[name]
    disabled = source.get("disabled_servers") or {}
    if name in disabled and isinstance(disabled[name], dict):
        return disabled[name]
    return None


def _is_mcp_disabled(source: dict, name: str) -> bool:
    if source["scope"] == "project":
        return name in (source.get("disabled_names") or set())
    return name in (source.get("disabled_servers") or {})


def _set_mcp_disabled(source: dict, name: str, disabled: bool) -> bool:
    if source["scope"] == "project":
        project = source["project"]
        disabled_list = project.get("disabledMcpjsonServers")
        if not isinstance(disabled_list, list):
            disabled_list = []
            project["disabledMcpjsonServers"] = disabled_list
        if disabled and name not in disabled_list:
            disabled_list.append(name)
        elif not disabled:
            project["disabledMcpjsonServers"] = [n for n in disabled_list if n != name]
        return True

    servers = source["servers"]
    disabled_servers = source["disabled_servers"]
    if disabled:
        cfg = servers.pop(name, disabled_servers.get(name))
        if cfg is not None:
            cfg.pop("disabled", None)
            disabled_servers[name] = cfg
    else:
        cfg = disabled_servers.pop(name, servers.get(name))
        if cfg is not None:
            cfg.pop("disabled", None)
            servers[name] = cfg
    return False


def _mask_mapping(values: Optional[dict]) -> dict:
    if not values:
        return {}
    masked = {}
    for k, v in values.items():
        value = str(v)
        if any(s in k.lower() for s in ("token", "key", "secret", "password", "credential", "auth")):
            masked[k] = value[:4] + "***" if len(value) > 4 else "***"
        else:
            masked[k] = value
    return masked


def _mcp_transport(cfg: dict) -> str:
    transport = str(cfg.get("type") or cfg.get("transport") or "").strip().lower()
    if transport:
        return transport
    if cfg.get("url"):
        return "http"
    return "stdio"


def _format_mcp_server(name: str, cfg: dict, source: dict, disabled: bool) -> dict:
    transport = _mcp_transport(cfg)
    return {
        "name": name,
        "scope": source["scope"],
        "scope_label": _source_label(source["scope"]),
        "config_path": str(source["path"]),
        "type": transport,
        "command": cfg.get("command", ""),
        "args": cfg.get("args", []) if isinstance(cfg.get("args", []), list) else [],
        "url": cfg.get("url", ""),
        "env": _mask_mapping(cfg.get("env") if isinstance(cfg.get("env"), dict) else {}),
        "headers": _mask_mapping(cfg.get("headers") if isinstance(cfg.get("headers"), dict) else {}),
        "disabled": disabled,
    }


def _stdio_config_from_request(req: "McpServerRequest") -> dict:
    command = (req.command or "").strip()
    if not command:
        raise HTTPException(status_code=400, detail="command is required")
    cfg = {
        "type": "stdio",
        "command": command,
        "args": req.args or [],
    }
    if req.env:
        cfg["env"] = req.env
    return cfg


class McpServerRequest(BaseModel):
    command: Optional[str] = None
    args: Optional[List[str]] = None
    env: Optional[Dict[str, str]] = None
    disabled: Optional[bool] = None


class McpServerPatchRequest(BaseModel):
    command: Optional[str] = None
    args: Optional[List[str]] = None
    env: Optional[Dict[str, str]] = None
    disabled: Optional[bool] = None


@app.get("/api/mcp/servers")
async def list_mcp_servers(cwd: Optional[str] = Query(default=None)):
    sources = _mcp_sources(cwd)
    result = []
    for source in sources:
        servers = source.get("servers") or {}
        for name, cfg in servers.items():
            if isinstance(cfg, dict):
                result.append(_format_mcp_server(name, cfg, source, _is_mcp_disabled(source, name)))
        for name, cfg in (source.get("disabled_servers") or {}).items():
            if name not in servers and isinstance(cfg, dict):
                result.append(_format_mcp_server(name, cfg, source, True))
    return {
        "servers": result,
        "cwd": str(_resolve_mcp_cwd(cwd)),
        "config_path": str(_CLAUDE_CONFIG_PATH),
        "config_paths": {
            "user": str(_CLAUDE_CONFIG_PATH),
            "local": str(_CLAUDE_CONFIG_PATH),
            "project": str(_resolve_mcp_cwd(cwd) / _PROJECT_MCP_FILENAME),
        },
    }


@app.post("/api/mcp/servers/{name}")
async def add_mcp_server(
    name: str,
    req: McpServerRequest,
    cwd: Optional[str] = Query(default=None),
    scope: str = Query(default="local"),
):
    target = _mcp_target(scope, cwd)
    if _mcp_config_in_source(target, name) is not None:
        raise HTTPException(status_code=409, detail=f"server '{name}' already exists")
    cfg = _stdio_config_from_request(req)
    if req.disabled:
        if target["scope"] == "project":
            target["servers"][name] = cfg
            save_choices = _set_mcp_disabled(target, name, True)
            _save_mcp_source(target, save_claude_choices=save_choices)
        else:
            target["disabled_servers"][name] = cfg
            _save_mcp_source(target)
    else:
        target["servers"][name] = cfg
        _save_mcp_source(target)
    return {"ok": True}


@app.patch("/api/mcp/servers/{name}")
async def patch_mcp_server(
    name: str,
    req: McpServerPatchRequest,
    cwd: Optional[str] = Query(default=None),
    scope: Optional[str] = Query(default=None),
):
    target = _find_mcp_source(name, scope, cwd)
    cfg = _mcp_config_in_source(target, name)
    if cfg is None:
        raise HTTPException(status_code=404, detail=f"server '{name}' not found")
    if req.command is not None:
        command = req.command.strip()
        if not command:
            raise HTTPException(status_code=400, detail="command is required")
        cfg["command"] = command
        cfg.setdefault("type", "stdio")
    if req.args is not None:
        cfg["args"] = req.args
    if req.env is not None:
        cfg["env"] = req.env
    save_choices = False
    cfg.pop("disabled", None)
    if req.disabled is not None:
        save_choices = _set_mcp_disabled(target, name, req.disabled)
    _save_mcp_source(target, save_claude_choices=save_choices)
    return {"ok": True}


@app.delete("/api/mcp/servers/{name}")
async def delete_mcp_server(
    name: str,
    cwd: Optional[str] = Query(default=None),
    scope: Optional[str] = Query(default=None),
):
    target = _find_mcp_source(name, scope, cwd)
    target["servers"].pop(name, None)
    if target["scope"] == "project":
        save_choices = _set_mcp_disabled(target, name, False)
        _save_mcp_source(target, save_claude_choices=save_choices)
    else:
        target["disabled_servers"].pop(name, None)
        _save_mcp_source(target)
    return {"ok": True}


@app.get("/api/cwds")
async def list_cwds():
    with db_connect() as conn:
        rows = conn.execute(
            "SELECT cwd, MAX(updated_at) AS last FROM sessions WHERE cwd <> '' GROUP BY cwd ORDER BY last DESC LIMIT 10"
        ).fetchall()
    return [r["cwd"] for r in rows]


@app.get("/api/tags")
async def list_tags():
    with db_connect() as conn:
        rows = conn.execute("SELECT tags FROM sessions WHERE tags <> '' AND archived = 0").fetchall()
    counts: Dict[str, int] = defaultdict(int)
    for r in rows:
        for t in (r["tags"] or "").split(","):
            t = t.strip()
            if t:
                counts[t] += 1
    return [{"name": k, "count": v} for k, v in sorted(counts.items(), key=lambda x: -x[1])]


@app.get("/api/stats")
async def stats():
    total_cost = 0.0
    total_duration = 0.0
    total_turns = 0
    daily: Dict[str, Dict[str, float]] = defaultdict(lambda: {"cost": 0.0, "turns": 0})
    tool_counts: Dict[str, int] = defaultdict(int)
    with db_connect() as conn:
        rows = conn.execute("SELECT id FROM sessions").fetchall()
    total_sessions = len(rows)
    for row in rows:
        events = load_events(row["id"])
        for ev in events:
            t = ev.get("type")
            if t == "result":
                cost = float(ev.get("total_cost_usd") or 0)
                dur = float(ev.get("duration_ms") or 0)
                ts = float(ev.get("ts") or time.time())
                total_cost += cost
                total_duration += dur
                total_turns += 1
                day = time.strftime("%Y-%m-%d", time.localtime(ts))
                daily[day]["cost"] += cost
                daily[day]["turns"] += 1
            elif t == "assistant":
                content = (ev.get("message") or {}).get("content") or []
                for block in content:
                    if block.get("type") == "tool_use":
                        tool_counts[block.get("name", "?")] += 1
    daily_sorted = sorted(daily.items(), key=lambda x: x[0])
    return {
        "total_cost_usd": round(total_cost, 4),
        "total_duration_ms": total_duration,
        "total_sessions": total_sessions,
        "total_turns": total_turns,
        "daily": [{"date": d, "cost": round(v["cost"], 4), "turns": v["turns"]} for d, v in daily_sorted],
        "tools": sorted(
            [{"name": k, "count": v} for k, v in tool_counts.items()],
            key=lambda x: -x["count"],
        )[:10],
    }


async def _list_files_via_git(base: Path, q_lower: str, limit: int) -> Optional[List[dict]]:
    try:
        proc = await asyncio.create_subprocess_exec(
            "git", "-C", str(base), "ls-files",
            "--cached", "--others", "--exclude-standard",
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.DEVNULL,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=5)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            return None
        if proc.returncode != 0:
            return None
    except (FileNotFoundError, asyncio.TimeoutError):
        return None
    except Exception:
        return None

    results: List[dict] = []
    for rel in stdout.decode("utf-8", errors="replace").splitlines():
        rel = rel.strip()
        if not rel:
            continue
        if q_lower and q_lower not in rel.lower():
            continue
        results.append({"path": str(base / rel), "rel": rel})
        if len(results) >= limit:
            break
    return results


@app.get("/api/files")
async def list_files(cwd: str = Query(...), q: str = Query(default=""), limit: int = Query(default=30)):
    base = Path(os.path.expanduser(cwd)).resolve()
    if not base.exists() or not base.is_dir():
        return []
    q_lower = q.lower()

    git_results = await _list_files_via_git(base, q_lower, limit)
    if git_results is not None:
        return git_results

    results: List[dict] = []
    for root, dirs, files in os.walk(str(base)):
        dirs[:] = [d for d in dirs if d not in IGNORED_DIRS and not d.startswith(".")]
        for f in files:
            if f.startswith("."):
                continue
            full = Path(root) / f
            try:
                rel = str(full.relative_to(base))
            except ValueError:
                continue
            if q_lower and q_lower not in rel.lower():
                continue
            results.append({"path": str(full), "rel": rel})
            if len(results) >= limit:
                return results
    return results


@app.get("/api/git")
async def git_status(cwd: str = Query(...)):
    target = os.path.expanduser(cwd)
    if not os.path.isdir(target):
        return {"branch": "", "dirty": 0, "available": False}
    try:
        proc = await asyncio.create_subprocess_exec(
            "git", "-C", target, "status", "--porcelain=v1", "--branch",
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=5)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            return {"branch": "", "dirty": 0, "available": False}
    except Exception:
        return {"branch": "", "dirty": 0, "available": False}
    if proc.returncode != 0:
        return {"branch": "", "dirty": 0, "available": False}
    branch = ""
    dirty = 0
    for line in stdout.decode("utf-8", errors="replace").splitlines():
        if line.startswith("##"):
            header = line[2:].strip()
            branch = header.split("...")[0].strip()
        else:
            dirty += 1
    return {"branch": branch, "dirty": dirty, "available": True}


app.mount("/uploads", StaticFiles(directory=str(UPLOADS_DIR)), name="uploads")


class _TextExtractor(HTMLParser):
    _SKIP_TAGS = {"script", "style", "noscript", "svg", "iframe", "head"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._parts: List[str] = []
        self._skip_depth = 0
        self._table_depth = 0
        self._active_rowspans: Dict[int, int] = {}
        self._new_rowspans: Dict[int, int] = {}
        self._current_row: Optional[List[str]] = None
        self._current_cell: Optional[List[str]] = None
        self._current_colspan = 1
        self._current_rowspan = 1

    def handle_starttag(self, tag: str, attrs: List) -> None:
        if tag in self._SKIP_TAGS:
            self._skip_depth += 1
        elif self._skip_depth > 0:
            return
        elif tag == "table":
            if self._table_depth == 0:
                self._parts.append("\n")
                self._active_rowspans = {}
            self._table_depth += 1
        elif self._table_depth > 0:
            if self._table_depth == 1 and tag == "tr":
                self._start_table_row()
            elif self._table_depth == 1 and tag in {"td", "th"}:
                self._start_table_cell(attrs)
            elif tag == "br" and self._current_cell is not None:
                self._current_cell.append("\n")
        elif tag in {"p", "br", "div", "li", "h1", "h2", "h3", "h4", "h5", "h6", "tr"}:
            self._parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP_TAGS and self._skip_depth > 0:
            self._skip_depth -= 1
        elif self._skip_depth > 0:
            return
        elif tag in {"td", "th"} and self._table_depth == 1:
            self._end_table_cell()
        elif tag == "tr" and self._table_depth == 1:
            self._end_table_row()
        elif tag == "table" and self._table_depth > 0:
            if self._table_depth == 1:
                self._end_table_cell()
                self._end_table_row()
                self._parts.append("\n")
            self._table_depth -= 1

    def handle_data(self, data: str) -> None:
        if self._skip_depth > 0:
            return
        chunk = data.strip()
        if chunk:
            if self._table_depth > 0 and self._current_cell is not None:
                self._current_cell.append(chunk)
            elif self._table_depth == 0:
                self._parts.append(chunk)

    def get_text(self) -> str:
        raw = " ".join(self._parts)
        collapsed = re.sub(r"[ \t]+", " ", raw)
        collapsed = re.sub(r"\n\s*", "\n", collapsed)
        collapsed = re.sub(r" +\n", "\n", collapsed)
        return re.sub(r"\n{3,}", "\n\n", collapsed).strip()

    def _span_value(self, attrs: List, name: str) -> int:
        value = dict(attrs).get(name)
        try:
            return max(1, min(int(value or 1), 100))
        except ValueError:
            return 1

    def _start_table_row(self) -> None:
        self._end_table_cell()
        self._end_table_row()
        self._current_row = []
        self._new_rowspans = {}

    def _start_table_cell(self, attrs: List) -> None:
        if self._current_row is None:
            self._start_table_row()
        self._end_table_cell()
        self._current_cell = []
        self._current_colspan = self._span_value(attrs, "colspan")
        self._current_rowspan = self._span_value(attrs, "rowspan")

    def _end_table_cell(self) -> None:
        if self._current_cell is None or self._current_row is None:
            return

        col = len(self._current_row)
        while self._active_rowspans.get(col, 0) > 0:
            self._current_row.append("")
            col += 1

        text = re.sub(r"\s+", " ", " ".join(self._current_cell)).strip()
        for offset in range(self._current_colspan):
            self._current_row.append(text if offset == 0 else "")
            if self._current_rowspan > 1:
                self._new_rowspans[col + offset] = max(
                    self._new_rowspans.get(col + offset, 0),
                    self._current_rowspan - 1,
                )

        self._current_cell = None
        self._current_colspan = 1
        self._current_rowspan = 1

    def _end_table_row(self) -> None:
        if self._current_row is None:
            return

        self._end_table_cell()
        if self._active_rowspans:
            max_col = max(self._active_rowspans)
            while len(self._current_row) <= max_col:
                self._current_row.append("")

        if any(cell for cell in self._current_row):
            self._parts.append("| " + " | ".join(self._current_row) + " |")
            self._parts.append("\n")

        next_rowspans = {
            col: remaining - 1
            for col, remaining in self._active_rowspans.items()
            if remaining > 1
        }
        for col, remaining in self._new_rowspans.items():
            next_rowspans[col] = max(next_rowspans.get(col, 0), remaining)

        self._active_rowspans = next_rowspans
        self._new_rowspans = {}
        self._current_row = None


def _extract_html_text(html: str) -> str:
    extractor = _TextExtractor()
    extractor.feed(html)
    extractor.close()
    return extractor.get_text()


def _is_private_host(host: str) -> bool:
    try:
        ip = ipaddress.ip_address(host)
        return ip.is_private or ip.is_loopback or ip.is_link_local or ip.is_reserved
    except ValueError:
        try:
            resolved = socket.gethostbyname(host)
            ip = ipaddress.ip_address(resolved)
            return ip.is_private or ip.is_loopback or ip.is_link_local or ip.is_reserved
        except Exception:
            return True


@app.post("/api/fetch-url")
async def fetch_url(req: FetchUrlRequest):
    parsed = urlparse(req.url)
    if parsed.scheme not in ("http", "https"):
        raise HTTPException(status_code=400, detail="only http/https allowed")
    if not parsed.hostname:
        raise HTTPException(status_code=400, detail="invalid url")
    if _is_private_host(parsed.hostname):
        raise HTTPException(status_code=400, detail="refusing to fetch private/internal host")

    def _do_fetch() -> Dict[str, str]:
        request = urllib.request.Request(
            req.url,
            headers={
                "User-Agent": "Mozilla/5.0 (compatible; ClaudeWeb/1.0)",
                "Accept": "text/html,application/xhtml+xml,text/plain;q=0.9,*/*;q=0.5",
            },
        )
        with urllib.request.urlopen(request, timeout=10) as resp:
            content_type = resp.headers.get("Content-Type", "") or ""
            charset = "utf-8"
            if "charset=" in content_type:
                charset = content_type.split("charset=", 1)[1].split(";")[0].strip()
            raw = resp.read(2 * 1024 * 1024)
        html = raw.decode(charset, errors="replace")
        title_match = re.search(r"<title[^>]*>(.*?)</title>", html, re.IGNORECASE | re.DOTALL)
        title = re.sub(r"\s+", " ", title_match.group(1)).strip() if title_match else req.url
        text = _extract_html_text(html)
        return {"title": title, "content": text}

    try:
        result = await asyncio.get_event_loop().run_in_executor(None, _do_fetch)
    except urllib.error.HTTPError as e:
        raise HTTPException(status_code=502, detail=f"remote {e.code}")
    except urllib.error.URLError as e:
        raise HTTPException(status_code=502, detail=f"fetch failed: {e.reason}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

    limit = max(500, min(req.max_chars or 10000, 50000))
    content = result["content"][:limit]
    return {
        "url": req.url,
        "title": result["title"] or req.url,
        "content": content,
        "truncated": len(result["content"]) > limit,
        "length": len(result["content"]),
    }


@app.get("/api/version")
async def get_version():
    return {"version": __version__}


@app.get("/changelog.json")
async def get_changelog():
    path = STATIC_DIR / "changelog.json"
    if not path.exists():
        return {"releases": []}
    return FileResponse(path, media_type="application/json")


@app.get("/")
async def index():
    return FileResponse(STATIC_DIR / "index.html")


def _check_claude_cli() -> Optional[str]:
    """Return claude CLI version string if available, else None."""
    import shutil
    import subprocess

    if shutil.which("claude") is None:
        return None
    try:
        result = subprocess.run(
            ["claude", "--version"],
            capture_output=True,
            text=True,
            timeout=5,
        )
        if result.returncode == 0:
            return result.stdout.strip() or "unknown"
    except (subprocess.TimeoutExpired, OSError):
        pass
    return "unknown"


def main():
    """CLI entry point for `claude-web` command."""
    import argparse
    import sys
    import uvicorn

    parser = argparse.ArgumentParser(description="Claude Code Web - Web UI for Claude Code CLI")
    parser.add_argument("--port", "-p", type=int, default=int(os.environ.get("PORT", "8765")), help="Port to listen on (default: 8765)")
    parser.add_argument("--host", type=str, default="127.0.0.1", help="Host to bind (default: 127.0.0.1)")
    parser.add_argument("--open", action="store_true", help="Open browser after starting")
    parser.add_argument("--version", "-v", action="store_true", help="Show version")
    parser.add_argument("--skip-cli-check", action="store_true", help="Skip claude CLI availability check on startup")
    args = parser.parse_args()

    if args.version:
        print(f"claude-web {__version__}")
        return

    print(f"Claude Code Web v{__version__}")
    print(f"  → http://{args.host}:{args.port}")
    print(f"  → Data: {_DATA_DIR}")

    if not args.skip_cli_check:
        claude_version = _check_claude_cli()
        if claude_version is None:
            print()
            print("  ✗ claude CLI not found in PATH", file=sys.stderr)
            print("    claude-web wraps the Claude Code CLI — install it first:", file=sys.stderr)
            print("      npm install -g @anthropic-ai/claude-code", file=sys.stderr)
            print("    Then run `claude` once to log in. Docs: https://docs.claude.com/claude-code", file=sys.stderr)
            print("    (Use --skip-cli-check to bypass this check.)", file=sys.stderr)
            print()
            sys.exit(1)
        print(f"  → Claude CLI: {claude_version}")

    _LOCAL_HOSTS = {"127.0.0.1", "localhost", "::1"}
    if args.host not in _LOCAL_HOSTS:
        print()
        print(f"  ⚠️  WARNING: binding to {args.host} exposes the server beyond localhost.", file=sys.stderr)
        print("     claude-web has NO built-in authentication. Anyone who can reach this", file=sys.stderr)
        print("     address can run commands, read your files, and burn your Claude quota.", file=sys.stderr)
        print("     Only use --host on a trusted network (e.g. tailscale, VPN, SSH tunnel).", file=sys.stderr)

    print()

    if args.open:
        import webbrowser
        import threading
        threading.Timer(1.5, lambda: webbrowser.open(f"http://{args.host}:{args.port}")).start()

    uvicorn.run(app, host=args.host, port=args.port)


if __name__ == "__main__":
    main()
